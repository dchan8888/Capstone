from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # Helper to add text slide
    def add_slide(title_text, content_points):
        # Layout 1 is typically 'Title and Content'
        slide_layout = prs.slide_layouts[1] 
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = title_text
        
        # Content
        tf = slide.placeholders[1].text_frame
        for point in content_points:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
            p.font.size = Pt(24)

    # Slide 1: Title Slide
    # Layout 0 is typically 'Title Slide'
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "🧢 Cappy: The Savvy Capybara"
    subtitle.text = "UK Financial Literacy Agent\nAgents for Good Track"

    # Slide 2: The Problem
    add_slide("The Problem: Financial Illiteracy", [
        "⚠️ UK students leave school knowing Pythagoras but not Overdrafts.",
        "⚠️ 'Free Money' myths lead to long-term debt and credit damage.",
        "⚠️ Banking jargon (APR, AER, ISA) is intimidating and exclusionary."
    ])

    # Slide 3: The Solution
    add_slide("The Solution: Cappy", [
        "✅ A Financial Concierge, not just a chatbot.",
        "✅ RAG-Augmented: Accurate, UK-specific financial data.",
        "✅ Multi-Agent Architecture: Separates advice generation from safety compliance.",
        "✅ Persona: 'Chill vibes' to lower financial anxiety."
    ])

    # Slide 4: Architecture
    add_slide("Under the Hood: Multi-Agent System", [
        "1. The Librarian (Gemini 2.5): Retrieves facts & drafts friendly advice.",
        "2. The Guardian (Gemini 2.5): Reviews for safety & compliance.",
        "3. Tool Use: ChromaDB for RAG (Retrieval Augmented Generation).",
        "4. Deployment: Streamlit Cloud."
    ])

    # Slide 5: Demo Placeholder
    # Layout 2 is typically 'Section Header'
    slide_layout = prs.slide_layouts[2] 
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "LIVE DEMO"
    
    # Save the presentation
    prs.save('Cappy_Pitch_Deck.pptx')
    print("Presentation saved as 'Cappy_Pitch_Deck.pptx'")

if __name__ == "__main__":
    try:
        import pptx
        create_presentation()
    except ImportError:
        print("Please install python-pptx first using: pip install python-pptx")